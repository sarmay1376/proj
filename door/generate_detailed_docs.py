import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from fpdf import FPDF

# --- CONFIGURATION ---
PROJECT_NAME = "Smart Door Lock System"
AUTHOR = "Doctor Gemini & Professional Engineering Team"

def add_bullet_points(content_shape, points):
    tf = content_shape.text_frame
    tf.word_wrap = True
    for point in points:
        p = tf.add_paragraph()
        p.text = point
        p.level = 0
        p.font.size = Pt(18)

def create_presentation():
    print("[*] Generating Ultra-Detailed PowerPoint...")
    prs = Presentation()
    
    # Slide 1: Title
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "PRO SECURE: " + PROJECT_NAME
    slide.placeholders[1].text = f"Hardware Mastery & Security Logic\n{AUTHOR}"

    # Slide 2: Stepper Wiring Mastery
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "5-Wire Stepper Wiring (28BYJ-48)"
    add_bullet_points(slide.placeholders[1], [
        "ULN2003 Driver Board is MANDATORY for current control.",
        "IN1 (Driver) -> Arduino Mega Pin 22",
        "IN2 (Driver) -> Arduino Mega Pin 24 (Sequence optimized)",
        "IN3 (Driver) -> Arduino Mega Pin 23 (Sequence optimized)",
        "IN4 (Driver) -> Arduino Mega Pin 25",
        "+ / - Pins: Connect to EXTERNAL 5V for maximum torque.",
        "Common Ground: Ensure External GND is tied to Mega GND."
    ])

    # Slide 3: Keypad Matrix Mapping
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "4x4 Keypad Matrix Decoded"
    add_bullet_points(slide.placeholders[1], [
        "Pin 1-4 (L to R): ROW Interface (2, 3, 4, A1)",
        "Pin 5-8 (L to R): COL Interface (6, 7, 8, 9)",
        "Row 4 is on A1 to prevent collision with RFID SPI.",
        "Internal Pull-ups: Keypad library handles noise rejection.",
        "# Key: Defined as soft-interrupt for PIN verification."
    ])

    # Slide 4: RFID & SPI Bus
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "RFID (MFRC522) High-Speed SPI"
    add_bullet_points(slide.placeholders[1], [
        "SS (SDA) -> Pin 53 | RST -> Pin 5",
        "SCK -> Pin 52 | MOSI -> Pin 51 | MISO -> Pin 50",
        "Power: Use 3.3V ONLY (Connecting to 5V will burn the IC).",
        "Logic Level: Mega pins are 5V, use 1k resistors if possible for long-term safety."
    ])

    prs.save("Smart_Door_Lock_Ultra_Manual.pptx")
    print("[+] Presentation saved: Smart_Door_Lock_Ultra_Manual.pptx")

class PDF(FPDF):
    def header(self):
        self.set_fill_color(0, 51, 102) 
        self.rect(0, 0, 210, 35, 'F')
        self.set_text_color(255, 255, 255)
        self.set_font('Helvetica', 'B', 22)
        self.cell(0, 15, "HARDWARE SCHEMATIC: " + PROJECT_NAME.upper(), 0, 1, 'C')
        self.ln(10)

def create_pdf():
    print("[*] Generating Ultra-Detailed PDF Manual...")
    pdf = PDF()
    pdf.add_page()
    
    pdf.set_text_color(0, 0, 0)
    pdf.set_font("Helvetica", 'B', 14)
    pdf.cell(0, 10, "1. STEPPER MOTOR (28BYJ-48 + ULN2003)", 0, 1, 'L')
    pdf.set_font("Courier", '', 10)
    pdf.multi_cell(0, 6, "DRIVE PIN | MEGA PIN | DESCRIPTION\n----------|----------|------------\nIN1       | 22       | Phase A\nIN2       | 24       | Phase C (Logic optimized)\nIN3       | 23       | Phase B (Logic optimized)\nIN4       | 25       | Phase D\nVCC       | EXT 5V   | Red wire (usually common)\nGND       | MEGA GND | Logic common")
    pdf.ln(5)

    pdf.set_font("Helvetica", 'B', 14)
    pdf.cell(0, 10, "2. 4x4 MATRIX KEYPAD (MEMBRANE)", 0, 1, 'L')
    pdf.set_font("Courier", '', 10)
    pdf.multi_cell(0, 6, "KEYPAD PIN | MEGA PIN | FUNCTION\n-----------|----------|---------\nPIN 1 (L)  | 2        | ROW 1\nPIN 2      | 3        | ROW 2\nPIN 3      | 4        | ROW 3\nPIN 4      | A1       | ROW 4 (Special Re-route)\nPIN 5      | 6        | COL 1\nPIN 6      | 7        | COL 2\nPIN 7      | 8        | COL 3\nPIN 8 (R)  | 9        | COL 4")
    pdf.ln(5)

    pdf.set_font("Helvetica", 'B', 14)
    pdf.cell(0, 10, "3. MFRC522 RFID SCANNER", 0, 1, 'L')
    pdf.set_font("Courier", '', 10)
    pdf.multi_cell(0, 6, "RFID PIN | MEGA PIN | NOTES\n---------|----------|------\nSDA (SS) | 53       | SPI Slave Select\nSCK      | 52       | Clock Line\nMOSI     | 51       | Data In\nMISO     | 50       | Data Out\nIRQ      | N/C      | Not Connected\nGND      | GND      | Ground\nRST      | 5        | Reset Line\n3.3V     | 3.3V     | DO NOT CONNECT TO 5V")
    
    pdf.output("Smart_Door_Lock_Detailed_Schematics.pdf")
    print("[+] PDF Manual saved: Smart_Door_Lock_Detailed_Schematics.pdf")

if __name__ == "__main__":
    create_presentation()
    create_pdf()

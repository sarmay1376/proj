import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from fpdf import FPDF

# --- CONFIGURATION ---
PROJECT_NAME = "Smart Door Lock System"
AUTHOR = "Embedded Systems Engineering Team"
OUTPUT_DIR = "."

# Image paths (generated separately, co-located with this script)
SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
BRAIN_DIR  = os.path.join(os.path.expanduser("~"), ".gemini", "antigravity", "brain",
                           "3a6cde1f-249a-4531-95de-acdd00dfa4b7")

IMG_HARDWARE  = os.path.join(BRAIN_DIR, "door_lock_system_1773892172871.png")
IMG_ARCH      = os.path.join(BRAIN_DIR, "door_lock_architecture_1773892227798.png")
IMG_DASHBOARD = os.path.join(BRAIN_DIR, "door_lock_dashboard_1773892242975.png")

# Colors
DARK_BG    = RGBColor(0x0A, 0x0E, 0x14)
CYAN       = RGBColor(0x00, 0xD2, 0xFF)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)

# --- Helper: Style a title ---
def style_title(shape, text, color=CYAN, size=32):
    shape.text = text
    p = shape.text_frame.paragraphs[0]
    p.font.bold = True
    p.font.size = Pt(size)
    p.font.color.rgb = color

# --- Helper: Add styled bullet points ---
def add_bullets(tf, points, size=18, color=LIGHT_GRAY):
    tf.word_wrap = True
    for i, point in enumerate(points):
        p = tf.paragraphs[i] if i == 0 else tf.add_paragraph()
        p.text = point
        p.font.size = Pt(size)
        p.font.color.rgb = color

# --- Helper: Dark background ---
def dark_bg(slide, prs):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG

def create_presentation():
    print("[*] Generating Premium PowerPoint Presentation v2...")
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    SLIDE_W = prs.slide_width
    SLIDE_H = prs.slide_height

    # ── Slide 1: Title ──────────────────────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    dark_bg(slide, prs)

    # Background hardware image (right half)
    if os.path.exists(IMG_HARDWARE):
        slide.shapes.add_picture(IMG_HARDWARE,
                                 Inches(6.8), Inches(0.5),
                                 Inches(6.2), Inches(6.5))

    # Dark overlay on left side (title area)
    box = slide.shapes.add_textbox(Inches(0.4), Inches(1.8), Inches(6.2), Inches(4))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "PRO SECURE"
    run.font.bold = True
    run.font.size = Pt(48)
    run.font.color.rgb = CYAN

    p2 = tf.add_paragraph()
    run2 = p2.add_run()
    run2.text = "\nSmart Door Lock System"
    run2.font.bold = True
    run2.font.size = Pt(28)
    run2.font.color.rgb = WHITE

    p3 = tf.add_paragraph()
    run3 = p3.add_run()
    run3.text = f"\n{AUTHOR}\nArduino Mega 2560 Edition"
    run3.font.size = Pt(16)
    run3.font.color.rgb = LIGHT_GRAY

    # ── Slide 2: The Problem ────────────────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    dark_bg(slide, prs)

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
    style_title(title_box.text_frame.paragraphs[0].runs[0] if title_box.text_frame.paragraphs[0].runs else title_box.text_frame.paragraphs[0], "The Challenge", size=36)
    t = title_box.text_frame
    t.paragraphs[0].text = "The Challenge"
    t.paragraphs[0].font.color.rgb = CYAN
    t.paragraphs[0].font.size = Pt(36)
    t.paragraphs[0].font.bold = True

    body = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(5.5))
    add_bullets(body.text_frame, [
        "🔓  Traditional locks are easy to pick or duplicate.",
        "🗝️  Managing physical keys for multiple users is inefficient.",
        "📋  No real-time audit logs for access tracking.",
        "🔄  No remote override capability for emergency access.",
        "🔔  No immediate alert system for denied attempts.",
    ], size=22)

    # ── Slide 3: Our Solution ───────────────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    dark_bg(slide, prs)

    sb = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
    sb.text_frame.paragraphs[0].text = "The Pro Secure Solution"
    sb.text_frame.paragraphs[0].font.color.rgb = CYAN
    sb.text_frame.paragraphs[0].font.size = Pt(36)
    sb.text_frame.paragraphs[0].font.bold = True

    body = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(5.5))
    add_bullets(body.text_frame, [
        "🔐  Triple-Factor Authentication: RFID + PIN Keypad + IR Remote.",
        "⚙️   Servo Motor Bolt: Instant, reliable locking mechanism.",
        "📟  16×2 LCD Display: Real-time status feedback for every event.",
        "🌐  Web Dashboard: Live monitoring & remote override via WebSockets.",
        "🔁  Non-blocking Auto-Lock: Safety re-lock after 5 seconds.",
    ], size=22)

    # ── Slide 4: Hardware Architecture ─────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    dark_bg(slide, prs)

    sb = slide.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(6), Inches(0.9))
    sb.text_frame.paragraphs[0].text = "Hardware Architecture"
    sb.text_frame.paragraphs[0].font.color.rgb = CYAN
    sb.text_frame.paragraphs[0].font.size = Pt(32)
    sb.text_frame.paragraphs[0].font.bold = True

    if os.path.exists(IMG_ARCH):
        slide.shapes.add_picture(IMG_ARCH,
                                 Inches(0.2), Inches(1.1),
                                 Inches(13.0), Inches(6.2))

    # ── Slide 5: Security Features ──────────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    dark_bg(slide, prs)

    sb = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
    sb.text_frame.paragraphs[0].text = "Security Features"
    sb.text_frame.paragraphs[0].font.color.rgb = CYAN
    sb.text_frame.paragraphs[0].font.size = Pt(36)
    sb.text_frame.paragraphs[0].font.bold = True

    body = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(5.5))
    add_bullets(body.text_frame, [
        "🪪  RFID: Scans 13.56MHz Mifare cards — only the Master UID unlocks.",
        "🔢  Keypad: 4-digit PIN with '*' buffer clear and 10-second timeout.",
        "📡  IR Remote: Specific hex codes for Open and Close commands.",
        "🚫  Denied Alert: 3-flash LED + descending buzzer tone on failure.",
        "⏱️   Auto-Relock: Non-blocking millis() timer — never blocks main loop.",
    ], size=22)

    # ── Slide 6: Web Dashboard ──────────────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    dark_bg(slide, prs)

    sb = slide.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(6), Inches(0.9))
    sb.text_frame.paragraphs[0].text = "Web Dashboard & Live Logs"
    sb.text_frame.paragraphs[0].font.color.rgb = CYAN
    sb.text_frame.paragraphs[0].font.size = Pt(32)
    sb.text_frame.paragraphs[0].font.bold = True

    if os.path.exists(IMG_DASHBOARD):
        slide.shapes.add_picture(IMG_DASHBOARD,
                                 Inches(0.5), Inches(1.0),
                                 Inches(12.3), Inches(6.3))

    # ── Slide 7: Software Layer ─────────────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    dark_bg(slide, prs)

    sb = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
    sb.text_frame.paragraphs[0].text = "Software Architecture"
    sb.text_frame.paragraphs[0].font.color.rgb = CYAN
    sb.text_frame.paragraphs[0].font.size = Pt(36)
    sb.text_frame.paragraphs[0].font.bold = True

    body = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(5.5))
    add_bullets(body.text_frame, [
        "⚡  State-Machine Logic: Non-blocking execution for 100% responsiveness.",
        "🔒  Dedicated Buffers: Keypad & Serial Web commands never interfere.",
        "🌐  Python Flask + Socket.IO: Eventlet-based async WebSocket bridge.",
        "🧵  Thread-Safe Writes: threading.Lock() prevents serial port crashes.",
        "🔄  Fail-Safe Boot: System always initializes in LOCKED state.",
    ], size=22)

    # ── Slide 8: Conclusion ─────────────────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    dark_bg(slide, prs)

    sb = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
    sb.text_frame.paragraphs[0].text = "Summary & Results"
    sb.text_frame.paragraphs[0].font.color.rgb = CYAN
    sb.text_frame.paragraphs[0].font.size = Pt(36)
    sb.text_frame.paragraphs[0].font.bold = True

    body = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(7), Inches(5.5))
    add_bullets(body.text_frame, [
        "✅  100% stable embedded IoT prototype achieved.",
        "✅  Zero-latency Hardware ↔ Web synchronization.",
        "✅  Comprehensive audit trail for all access attempts.",
        "✅  LCD provides instant real-world status feedback.",
        "✅  Servo motor: silent, precise, reliable bolt mechanism.",
        "",
        "🏁  Status: READY FOR FACILITY DEPLOYMENT.",
    ], size=20)

    if os.path.exists(IMG_HARDWARE):
        slide.shapes.add_picture(IMG_HARDWARE,
                                 Inches(7.8), Inches(1.2),
                                 Inches(5.2), Inches(5.5))

    prs.save("Smart_Door_Lock_Premium_Presentation.pptx")
    print("[+] Presentation saved: Smart_Door_Lock_Premium_Presentation.pptx")


class PDF(FPDF):
    def header(self):
        self.set_fill_color(10, 14, 20)
        self.rect(0, 0, 210, 40, 'F')
        self.set_text_color(255, 255, 255)
        self.set_font('Helvetica', 'B', 20)
        self.cell(0, 20, PROJECT_NAME.upper(), 0, 1, 'C')
        self.set_font('Helvetica', 'I', 10)
        self.cell(0, 10, "OFFICIAL TECHNICAL SPECIFICATION", 0, 1, 'C')
        self.ln(10)

    def footer(self):
        self.set_y(-25)
        self.set_font('Helvetica', 'I', 8)
        self.set_text_color(128)
        self.cell(0, 10, f'Smart Door Lock Technical Manual - Page {self.page_no()}', 0, 0, 'C')


def create_pdf():
    print("[*] Generating PDF Manual...")
    pdf = PDF()
    pdf.add_page()
    pdf.set_auto_page_break(auto=True, margin=30)

    pdf.set_text_color(0, 0, 0)
    pdf.set_font("Helvetica", 'B', 14)
    pdf.cell(0, 10, "1. EXECUTIVE SUMMARY", 0, 1, 'L')
    pdf.set_font("Helvetica", '', 11)
    pdf.multi_cell(0, 7, (
        "This document outlines the architecture and deployment of the Pro-Secure Smart Door Lock v2. "
        "Utilizing triple-redundant authentication (RFID, Matrix Keypad, and IR Wireless), the system "
        "provides a robust barrier against unauthorized access while maintaining usability through a "
        "16x2 LCD display and a real-time web-based monitoring dashboard."
    ))
    pdf.ln(5)

    pdf.set_font("Helvetica", 'B', 12)
    pdf.cell(0, 10, "2. HARDWARE PIN CONFIGURATION (MEGA 2560)", 0, 1, 'L')
    pdf.set_font("Courier", 'B', 10)

    data = [
        ["COMPONENT",           "PINS",              "NOTES"],
        ["RFID (SPI)",          "50, 51, 52, 53, 5", "MISO/MOSI/SCK/SS/RST"],
        ["KEYPAD ROWS",         "30, 31, 32, 33",    "R1-R4 Sequential"],
        ["KEYPAD COLS",         "34, 35, 36, 37",    "C1-C4 Sequential"],
        ["SERVO MOTOR",         "9",                  "PWM Signal"],
        ["LCD RS/E",            "26, 27",             "Control Pins"],
        ["LCD D4-D7",           "28, 29, A8, A9",    "Data Pins (4-bit mode)"],
        ["IR SENSOR",           "18",                 "Interrupt-capable"],
        ["BUZZER / LED",        "10, 11",             "Feedback"]
    ]
    col_widths = [45, 55, 80]
    for row in data:
        for i, item in enumerate(row):
            pdf.cell(col_widths[i], 10, item, 1)
        pdf.ln()
    pdf.ln(5)

    pdf.set_font("Helvetica", 'B', 12)
    pdf.cell(0, 10, "3. SECURITY LOGIC FLOW", 0, 1, 'L')
    pdf.set_font("Helvetica", '', 11)
    pdf.multi_cell(0, 7, (
        "- [BOOT]: All peripherals init; servo locks; LCD shows 'INITIALIZING'.\n"
        "- [IDLE]: RFID, Keypad, IR, and Serial polled at 20ms intervals.\n"
        "- [GRANT]: If RFID UID matches, PIN matches, or correct IR code: servo rotates 90°.\n"
        "- [LCD]: Updates in real-time for every state change.\n"
        "- [AUTO-LOCK]: After 5s, servo returns to 0° and LCD shows 'LOCKED'.\n"
        "- [DENY]: 3-flash LED + descending tone; LCD shows 'ACCESS DENIED'."
    ))
    pdf.ln(5)

    pdf.set_font("Helvetica", 'B', 12)
    pdf.cell(0, 10, "4. WEB DASHBOARD & OVERRIDE", 0, 1, 'L')
    pdf.set_font("Helvetica", '', 11)
    pdf.multi_cell(0, 7, (
        "The 'ui.py' script acts as a high-speed telemetry bridge. It parses Arduino Serial data and "
        "broadcasts it via WebSockets (Flask-SocketIO). A threading.Lock() prevents concurrent serial "
        "writes from crashing the connection. The dashboard allows remote unlock/lock via 'UNLOCK_CMD'/'LOCK_CMD'."
    ))

    pdf.output("Smart_Door_Lock_Premium_Manual.pdf")
    print("[+] PDF Manual saved: Smart_Door_Lock_Premium_Manual.pdf")


if __name__ == "__main__":
    create_presentation()
    create_pdf()

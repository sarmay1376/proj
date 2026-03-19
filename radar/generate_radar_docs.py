import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from fpdf import FPDF

# --- CONFIG ---
PROJECT_NAME = "RADAR PRO: Advanced Detection System"
AUTHOR = "Doctor Gemini & Professional Engineering Team"

def add_bullet_points(content_shape, points):
    tf = content_shape.text_frame
    tf.word_wrap = True
    for point in points:
        p = tf.add_paragraph()
        p.text = point
        p.level = 0
        p.font.size = Pt(18)

def create_presentation():
    print("[*] Generating Ultra-Detailed Radar PowerPoint...")
    prs = Presentation()
    
    # Slide 1: Title
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "PRO RADAR: " + PROJECT_NAME
    slide.placeholders[1].text = f"Ultrasonic Sweep & Real-time Visualization\n{AUTHOR}"

    # Slide 2: Hardware Wiring (The "Which Wire" slide)
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Precision Wiring Schematic"
    add_bullet_points(slide.placeholders[1], [
        "HC-SR04: Trig=10, Echo=11 (5V Supply)",
        "Servo: Signal=9 (PWM), VCC=5V, GND=GND",
        "LCD 16x2: RS=2, E=3, D4-D7=4, 5, 6, 7",
        "Buzzer: Pin 8 | Red LED: Pin A5",
        "Final optimized professional layout."
    ])

    # Slide 3: Embedded Fritzing
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Visual Hardware Layout"
    if os.path.exists("wiring_diagram.png"):
        slide.shapes.add_picture("wiring_diagram.png", Inches(1), Inches(1.5), width=Inches(8))

    # Slide 4: Algorithm & Filtering
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Detection Algorithm"
    add_bullet_points(slide.placeholders[1], [
        "Speed of Sound Math: (duration * 0.034) / 2",
        "Alpha Filter: 70% New Data / 30% Previous History",
        "Auto-Alert: Any valid echo (1-400cm) triggers buzzer.",
        "Smooth Sweep: 30ms delay per degree for stabilization."
    ])

    prs.save("Radar_Pro_Edition_Presentation.pptx")

class PDF(FPDF):
    def header(self):
        self.set_fill_color(0, 51, 102) 
        self.rect(0, 0, 210, 35, 'F')
        self.set_text_color(255, 255, 255)
        self.set_font('Helvetica', 'B', 22)
        self.cell(0, 15, "HARDWARE SCHEMATIC: RADAR PRO", 0, 1, 'C')
        self.ln(10)

def create_pdf():
    print("[*] Generating Ultra-Detailed Radar PDF Manual...")
    pdf = PDF()
    pdf.add_page()
    
    pdf.set_text_color(0, 0, 0)
    pdf.set_font("Helvetica", 'B', 14)
    pdf.cell(0, 10, "1. RADAR WIRING PINOUT", 0, 1, 'L')
    pdf.set_font("Courier", '', 10)
    pdf.multi_cell(0, 6, "COMPONENT | ARDUINO PIN | NOTES\n----------|-------------|------------\nHC-SR04 T | 10          | Trigger\nHC-SR04 E | 11          | Echo pulse\nSERVO SIG | 9           | PWM\nBUZZER    | 8           | Active\nRED LED   | A5          | Alert\nLCD RS    | 2           | \nLCD E     | 3           | \nLCD D4-D7 | 4, 5, 6, 7  | Data bus")
    pdf.ln(5)

    pdf.set_font("Helvetica", 'B', 14)
    pdf.cell(0, 10, "2. CORE MEASUREMENT LOGIC", 0, 1, 'L')
    pdf.set_font("Courier", '', 9)
    pdf.multi_cell(0, 5, "duration = pulseIn(echoPin, HIGH);\ndistance = duration * 0.034 / 2;\n\n// Alpha filter in Python UI:\ndistances[idx] = (0.7 * current_dist) + (0.3 * distances[idx]);")
    
    if os.path.exists("wiring_diagram.png"):
        pdf.add_page()
        pdf.image("wiring_diagram.png", x=10, y=30, w=190)

    pdf.output("Radar_Pro_Technical_Manual.pdf")

if __name__ == "__main__":
    create_presentation()
    create_pdf()

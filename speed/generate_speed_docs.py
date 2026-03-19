import os
from pptx import Presentation
from pptx.util import Inches, Pt
from fpdf import FPDF

# --- CONFIG ---
PROJECT_NAME = "SPEED TRAP PRO: High-Velocity Enforcement"
AUTHOR = "Doctor Gemini & Professional Engineering Team"

def add_bullet_points(content_shape, points):
    tf = content_shape.text_frame
    tf.word_wrap = True
    for point in points:
        p = tf.add_paragraph()
        p.text = point
        p.level = 0
        p.font.size = Pt(18)

def create_presentation():
    print("[*] Generating Ultra-Detailed Speed Trap PowerPoint...")
    prs = Presentation()
    
    # Slide 1: Title
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "SPEED PRO: " + PROJECT_NAME
    slide.placeholders[1].text = f"Dual-Gate Velocity Logic & Web Dashboard\n{AUTHOR}"

    # Slide 2: Hardware Wiring
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Gate & Logic Wiring"
    add_bullet_points(slide.placeholders[1], [
        "Sensor 1 (Entry): Pins 3, 4",
        "Sensor 2 (Exit): Pins 5, 6",
        "Gap Distance: Flexible (Calibrated to 13cm)",
        "LCD: 7, 8, 9, 10, 11, 12 (Parallel 4-bit)",
        "Alert: Buzzer=A0, Red LED=A1"
    ])

    # Slide 3: Visual Schematic
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Hardware Interconnect"
    if os.path.exists("wiring_diagram.png"):
        slide.shapes.add_picture("wiring_diagram.png", Inches(1), Inches(1.5), width=Inches(8))

    # Slide 4: Speed Calculation Logic
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Math & Logic Deep Dive"
    add_bullet_points(slide.placeholders[1], [
        "Speed = 0.2 meters / (t2 - t1) in seconds",
        "Microsecond precision using micros()",
        "Double-Confirmation logic to prevent false starts",
        "Auto-Relock/Wait timeout after 4 seconds",
        "Threshold Alert: Fires if KM/H > 5.0"
    ])

    prs.save("Speed_Trap_Pro_Edition_Presentation.pptx")

class PDF(FPDF):
    def header(self):
        self.set_fill_color(102, 0, 0) # Dark red for speed alert
        self.rect(0, 0, 210, 35, 'F')
        self.set_text_color(255, 255, 255)
        self.set_font('Helvetica', 'B', 22)
        self.cell(0, 15, "HARDWARE SCHEMATIC: SPEED TRAP PRO", 0, 1, 'C')
        self.ln(10)

def create_pdf():
    print("[*] Generating Ultra-Detailed Speed Trap PDF Manual...")
    pdf = PDF()
    pdf.add_page()
    
    pdf.set_text_color(0, 0, 0)
    pdf.set_font("Helvetica", 'B', 14)
    pdf.cell(0, 10, "1. SPEED TRAP PINOUT", 0, 1, 'L')
    pdf.set_font("Courier", '', 10)
    pdf.multi_cell(0, 6, "COMPONENT | ARDUINO PIN | NOTES\n----------|-------------|------------\nSENSOR 1  | 3, 4        | Entry Gate\nSENSOR 2  | 5, 6        | Exit Gate\nBUZZER    | A0          | Speed Alert\nLED       | A1          | Speed Alert\nLCD       | 7, 8, 9, 10, 11, 12 | Status\nDISTANCE  | 13cm (Variable)")
    pdf.ln(5)

    pdf.set_font("Helvetica", 'B', 14)
    pdf.cell(0, 10, "2. CORE SPEED FORMULA", 0, 1, 'L')
    pdf.set_font("Courier", '', 10)
    pdf.multi_cell(0, 6, "float dt = (t2 - t1) / 1000000.0;\nfloat speedMS = DIST_M / dt;\nfloat speedKMH = speedMS * 3.6;")
    
    if os.path.exists("wiring_diagram.png"):
        pdf.add_page()
        pdf.image("wiring_diagram.png", x=10, y=30, w=190)

    pdf.output("Speed_Trap_Pro_Technical_Manual.pdf")

if __name__ == "__main__":
    create_presentation()
    create_pdf()
